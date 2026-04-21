import requests
import io
import re
import pandas as pd
from pathlib import Path
from typing import Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Paths & Branding ---
BASE_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = BASE_DIR.parent.parent

# REQUIRED CONSTANTS
DEFAULT_CSV_PATH = BASE_DIR / "data" / "vital_products.csv"
LOGO_PATH = PROJECT_ROOT / "static" / "logo_vital.png"

OLLAMA_URL = "http://localhost:11434/api/generate"
OLLAMA_MODEL = "llama3.2:latest"

COLOR_VITAL_BLUE = RGBColor(0, 51, 102)
COLOR_VITAL_TEAL = RGBColor(0, 153, 153)
COLOR_TEXT_DARK = RGBColor(33, 33, 33)

def clean_text(text: str) -> str:
    """Removes Markdown bolding stars and excessive whitespace."""
    if not text:
        return ""
    text = text.replace("**", "")
    text = text.replace("* ", "• ")
    text = re.sub(r'\*+', '', text)
    return text.strip()

def ask_ollama(prompt: str) -> str:
    try:
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.4,
                    "num_predict": 1000
                }
            },
            timeout=180
        )
        response.raise_for_status()
        return response.json().get("response", "").strip()
    except Exception as e:
        print(f"Ollama Error: {e}")
        return ""

def build_prompt(row: pd.Series) -> str:
    return f"""
Tu es le Stratège Marketing en Chef des Laboratoires Vital. Rédige un dossier de présentation haut de gamme.
PRODUIT : {row.get('name', '')}
INFOS : {row.get('indications', '')} | {row.get('forme', '')}

STRICT : Ne pas utiliser de formattage Markdown comme les étoiles (**).
Structure ta réponse avec ces balises :

[ACCROCHE]
Un titre prestigieux (max 15 mots).

[PROBLEMATIQUE]
Analyse riche du besoin patient (80 mots environ).

[SOLUTION_EXPERTE]
Explique le mode d'action unique et la supériorité Vital (100 mots environ).

[ARGUMENTS_CLES]
4 piliers stratégiques. Pour chaque pilier : Un titre suivi d'une phrase d'impact.
(Total max 150 mots).

[QUALITE_TECHNIQUE]
Détaille l'excellence de fabrication et la forme {row.get('forme', '')} (60 mots).

[CTA]
Conclusion de leadership et appel à l'action court.
"""

def parse_sections(content: str):
    def extract(tag):
        pattern = rf"\[{tag}\](.*?)(?=\[|$)"
        match = re.search(pattern, content, re.DOTALL | re.IGNORECASE)
        if match:
            return clean_text(match.group(1))
        return ""

    return {
        "accroche": extract("ACCROCHE"),
        "probleme": extract("PROBLEMATIQUE"),
        "solution": extract("SOLUTION_EXPERTE"),
        "arguments": extract("ARGUMENTS_CLES"),
        "technique": extract("QUALITE_TECHNIQUE"),
        "cta": extract("CTA")
    }

def _narrations_from_data(data: dict, product_name: str) -> list[str]:
    """
    Return one TTS narration string per slide (6 slides).
    Used by video_generation.py to drive the avatar's speech.
    """
    name = product_name
    return [
        # Slide 1 — Cover
        f"Bienvenue. Voici la présentation de {name}. {data['accroche']}",
        # Slide 2 — Problématique
        data["probleme"] or f"Découvrons ensemble les enjeux liés à {name}.",
        # Slide 3 — Solution
        data["solution"] or f"{name} apporte une réponse experte et innovante.",
        # Slide 4 — Arguments clés
        data["arguments"] or f"Voici les piliers stratégiques de {name}.",
        # Slide 5 — Qualité technique
        data["technique"] or f"L'excellence de fabrication de {name} est notre engagement.",
        # Slide 6 — Conclusion
        data["cta"] or f"Merci pour votre attention. {name} est disponible dès maintenant.",
    ]

def add_branding(slide, title_text: str = ""):
    """Sidebar, Smaller Logo, and Auto-fitting Header."""
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = COLOR_VITAL_TEAL
    sidebar.line.visible = False

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(8.8), Inches(0.2), width=Inches(0.8))

    if title_text:
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(7.8), Inches(1))
        tf = title_box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_VITAL_BLUE

def setup_text_frame(shape):
    """Utility to enable word wrap and auto-size for text boxes."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return tf

def create_presentation(row: pd.Series, ai_content: str, output_file: Path):
    prs = Presentation()
    data = parse_sections(ai_content)
    img_data = download_image(row.get('image_url', ''))

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(4.4), Inches(0.8), width=Inches(1.2))

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), Inches(10), Inches(3))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_VITAL_BLUE
    bg.line.visible = False

    name_box = slide.shapes.add_textbox(Inches(0), Inches(3.3), Inches(10), Inches(1))
    tf = setup_text_frame(name_box)
    p = tf.paragraphs[0]
    p.text = row['name']
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.5), Inches(10), Inches(1))
    tf = setup_text_frame(sub_box)
    p = tf.paragraphs[0]
    p.text = data['accroche'] or "L'innovation au service de votre santé."
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = COLOR_VITAL_TEAL
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Problématique ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Analyse et Enjeux")
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    tf = setup_text_frame(body)
    p = tf.paragraphs[0]
    p.text = data['probleme']
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_TEXT_DARK

    # --- Slide 3: Solution (Split with Image) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "La Réponse Vital Labs")
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.8), Inches(5))
    tf = setup_text_frame(tx_box)
    p = tf.paragraphs[0]
    p.text = data['solution']
    p.font.size = Pt(18)
    if img_data:
        slide.shapes.add_picture(img_data, Inches(5.6), Inches(1.5), height=Inches(4.8))

    # --- Slide 4: Arguments ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Piliers de Performance")
    arg_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = setup_text_frame(arg_box)
    p = tf.paragraphs[0]
    p.text = data['arguments']
    p.font.size = Pt(18)

    # --- Slide 5: Technique & Fiche Technique ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Engagement Qualité & Science")

    tech_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(2.2))
    tf_tech = setup_text_frame(tech_box)
    p_tech = tf_tech.paragraphs[0]
    p_tech.text = data['technique']
    p_tech.font.size = Pt(20)

    spec_map = [
        ("Catégorie", "categories"),
        ("Classe", "classe"),
        ("Forme", "forme"),
        ("Conditionnement", "infos_produit")
    ]

    valid_specs = []
    for label, key in spec_map:
        val = str(row.get(key, "")).strip()
        if val and val.lower() not in ["aucun", "aucune", "n/a", "nan", "none"]:
            valid_specs.append((label, val))

    if valid_specs:
        rows_n = len(valid_specs)
        table_shape = slide.shapes.add_table(
            rows_n, 2, Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.5 * rows_n)
        )
        table = table_shape.table
        table.columns[0].width = Inches(2.2)
        table.columns[1].width = Inches(6.2)

        for i, (label, value) in enumerate(valid_specs):
            cell_label = table.cell(i, 0)
            cell_label.fill.solid()
            cell_label.fill.fore_color.rgb = RGBColor(245, 245, 245)
            p_l = cell_label.text_frame.paragraphs[0]
            p_l.text = label
            p_l.font.bold = True
            p_l.font.size = Pt(18)
            p_l.font.color.rgb = COLOR_VITAL_BLUE

            cell_val = table.cell(i, 1)
            p_v = cell_val.text_frame.paragraphs[0]
            p_v.text = value
            p_v.font.size = Pt(18)
            p_v.font.color.rgb = COLOR_TEXT_DARK

    # --- Slide 6: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Synthèse & Engagement")
    cta_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    tf = setup_text_frame(cta_box)
    p = tf.paragraphs[0]
    p.text = data['cta']
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_VITAL_TEAL
    p.alignment = PP_ALIGN.CENTER

    prs.save(output_file)

def download_image(url: str) -> Optional[io.BytesIO]:
    if not url or not isinstance(url, str) or not url.startswith("http"):
        return None
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return io.BytesIO(response.content)
    except Exception:
        return None

# ─────────────────────────────────────────────────────────────────────────────
# Core generation helpers
# ─────────────────────────────────────────────────────────────────────────────

def _load_product_row(product_name: str, csv_path: Optional[Path]) -> pd.Series:
    """Load the CSV and return the row for the given product name (raises on miss)."""
    csv_file = Path(csv_path) if csv_path else DEFAULT_CSV_PATH
    if not csv_file.exists():
        raise FileNotFoundError(f"CSV file not found: {csv_file}")
    df = pd.read_csv(csv_file).fillna("")
    matched = df[df["name"].str.lower() == product_name.lower()]
    if matched.empty:
        raise ValueError(f"Produit '{product_name}' non trouvé.")
    return matched.iloc[0]


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def generate_presentation_for_product(
    product_name: str,
    csv_path: Optional[Path] = None,
    output_dir: Optional[Path] = None,
) -> Path:
    """Generate and save the PPTX.  Returns the output path."""
    pptx, _ = generate_presentation_for_product_with_narrations(
        product_name, csv_path, output_dir
    )
    return pptx


def generate_presentation_for_product_with_narrations(
    product_name: str,
    csv_path: Optional[Path] = None,
    output_dir: Optional[Path] = None,
) -> Tuple[Path, list[str]]:
    """
    Generate and save the PPTX, and return per-slide TTS narration strings.

    Returns:
        (pptx_path, slide_narrations)   — used by video_generation.py
    """
    row = _load_product_row(product_name, csv_path)

    ai_raw = ask_ollama(build_prompt(row))
    data   = parse_sections(ai_raw)

    out_dir = Path(output_dir) if output_dir else BASE_DIR / "generated_presentations"
    out_dir.mkdir(parents=True, exist_ok=True)

    output_file = out_dir / f"{product_name.replace(' ', '_')}.pptx"
    create_presentation(row, ai_raw, output_file)

    narrations = _narrations_from_data(data, product_name)
    return output_file, narrations